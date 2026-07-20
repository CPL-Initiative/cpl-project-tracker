#!/usr/bin/env python3
"""Fill the CPL section (slides 17-18 + one added) of the 2026 Annual Budget
Workshop template with the 3 CBO funding slides — native title style, CCC brand
palette (002F6D / 0066BA / FFB600 / 40B4E5 / 555759), Source Sans Pro."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SRC = "template.pptx"
OUT = "20260723_2026_Annual_Budget_Workshop_CPL.pptx"

NAVY = RGBColor(0x00,0x2F,0x6D)
BLUE = RGBColor(0x00,0x66,0xBA)
GOLD = RGBColor(0xFF,0xB6,0x00)
CYAN = RGBColor(0x40,0xB4,0xE5)
GREY = RGBColor(0x55,0x57,0x59)
DARK = RGBColor(0x33,0x38,0x3D)
WHITE= RGBColor(0xFF,0xFF,0xFF)
TINT = RGBColor(0xEE,0xF2,0xF8)   # light card tint
TINTB= RGBColor(0xE3,0xEC,0xF6)
LINE = RGBColor(0xD3,0xDD,0xEA)
FONT = "Source Sans Pro"

prs = Presentation(SRC)

def para(tf, runs, align=PP_ALIGN.LEFT, sa=4, sb=0, ls=1.0, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after=Pt(sa); p.space_before=Pt(sb); p.line_spacing=ls
    for (t,sz,col,b,*rest) in runs:
        it = rest[0] if rest else False
        r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.italic=it
        r.font.color.rgb=col; r.font.name=FONT
    return p

def textbox(slide,x,y,w,h,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    return tf

def box(slide,x,y,w,h,fill=None,ln=None,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp=slide.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb=ln; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    try: sp.adjustments[0]=0.06
    except: pass
    return sp

def set_title(slide, title, size=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx==0:
            tf=ph.text_frame; tf.clear()
            r=tf.paragraphs[0].add_run(); r.text=title
            r.font.name=FONT; r.font.bold=True; r.font.color.rgb=NAVY
            if size: r.font.size=Pt(size)
            return
def drop_body(slide):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx==11:
            ph._element.getparent().remove(ph._element)

def dot_list(slide,x,y,w,header,items,dot):
    tf=textbox(slide,x,y,w,0.4)
    para(tf,[(header,16,NAVY,True)],first=True)
    tb=textbox(slide,x,y+0.44,w,3.0)
    for i,it in enumerate(items):
        para(tb,[("●  ",12.5,dot,False),(it,13.5,DARK,False)],sa=7,ls=1.02,first=(i==0))

# ---------- SLIDE 1 (fill existing slide 17, index 16) ----------
s1 = prs.slides[16]; drop_body(s1)
set_title(s1,"Standing Up CPL at Every College")
tf=textbox(s1,0.6,1.62,12.13,0.72)
para(tf,[("SB 135 / AB 135 establishes CPL as a ",13.5,GREY,False),
         ("systemwide initiative",13.5,NAVY,True),
         (" under the Master Plan for Career Education. The one-time investment helps each college build local capacity to meet new duties.",13.5,GREY,False)],
     first=True,ls=1.1)
box(s1,0.6,2.42,5.86,2.85,fill=TINT,ln=LINE)
box(s1,6.87,2.42,5.86,2.85,fill=TINT,ln=LINE)
dot_list(s1,0.95,2.68,5.2,"What colleges now do",[
 "Evaluate incoming students’ prior learning & credentials",
 "Adopt statewide faculty credit recommendations",
 "Accept transcribed CPL from other colleges — credit that travels",
 "Build local capacity: people, processes & technology"], NAVY)
dot_list(s1,7.22,2.68,5.2,"What the funds seed",[
 "Dedicated CPL staffing & faculty engagement",
 "Student identification, outreach & advising",
 "Employer & industry partnerships",
 "New & stackable pathways — local CPL innovation"], GOLD)
box(s1,0.6,5.42,12.13,0.92,fill=NAVY)
tf=textbox(s1,0.95,5.42,11.45,0.92,anchor=MSO_ANCHOR.MIDDLE)
para(tf,[("One-time funds build ",13.5,WHITE,False),("local capacity",13.5,GOLD,True),
         ("; ongoing funds sustain the ",13.5,WHITE,False),("shared statewide backbone",13.5,GOLD,True),
         (" — the technology and faculty credit recommendations every college draws on.",13.5,WHITE,False)],
     first=True,ls=1.08)

# ---------- SLIDE 2 (fill existing slide 18, index 17) ----------
s2 = prs.slides[17]; drop_body(s2)
set_title(s2,"Three Funding Priorities")
tf=textbox(s2,0.6,1.66,12.13,0.7)
para(tf,[("Each college’s share is earned on measured CPL outcomes — reflecting the statute’s goal of advancing career attainment through CPL.",13.5,GREY,False,True)],first=True,ls=1.1)
pris=[("1","Completion","Grow certificate & degree completion — and career attainment — through CPL awards."),
      ("2","Access","Reach and enroll more students through CPL, and make it visible and accessible."),
      ("3","Capacity & Mobility","Build durable CPL — documentable, interoperable, and portable across colleges and segments.")]
x0,y0,w,h,g=0.6,2.55,3.92,2.55,0.185
for i,(n,name,desc) in enumerate(pris):
    x=x0+i*(w+g)
    box(s2,x,y0,w,h,fill=TINT,ln=LINE)
    box(s2,x+0.3,y0+0.28,0.62,0.62,fill=NAVY,shape=MSO_SHAPE.OVAL)
    tf=textbox(s2,x+0.3,y0+0.28,0.62,0.62,anchor=MSO_ANCHOR.MIDDLE)
    para(tf,[(n,22,WHITE,True)],align=PP_ALIGN.CENTER,first=True)
    tf=textbox(s2,x+0.3,y0+1.02,w-0.55,0.5)
    para(tf,[("PRIORITY "+n,10.5,BLUE,True)],first=True)
    para(tf,[(name,19,NAVY,True)],sb=1)
    tf=textbox(s2,x+0.3,y0+1.72,w-0.6,0.75)
    para(tf,[(desc,12.5,DARK,False)],first=True,ls=1.05)
box(s2,0.6,5.35,12.13,0.92,fill=NAVY)
tf=textbox(s2,0.95,5.35,11.45,0.92,anchor=MSO_ANCHOR.MIDDLE)
para(tf,[("DRAFT   ",13,GOLD,True),
         ("Priority weighting and dollar allocations are still in development with the field — the three-priority framework is what’s settled.",13,WHITE,False)],
     first=True,ls=1.06)

# ---------- SLIDE 3 (add new slide with same layout, then reorder) ----------
content_layout = prs.slides[16].slide_layout
s3 = prs.slides.add_slide(content_layout)
drop_body(s3)
set_title(s3,"Guiding Principles")
tf=textbox(s3,0.6,1.72,12.13,0.55)
para(tf,[("The fiscal design behind the model — how funds are allocated and safeguarded.",13.5,GREY,False,True)],first=True)
princ=[("Outcomes-based","Funding follows measured results — not headcount alone."),
       ("Equitable & inclusive","A minimum-viable floor so every participating college can build a real program; a rural-college allowance; noncredit-feeder support."),
       ("Transparent & predictable","One clear, published model — colleges can see how their allocation is derived."),
       ("Sustainable","One-time funds build local capacity; ongoing funds sustain the statewide backbone."),
       ("Systemwide & portable","Credit earned once travels — aligned across colleges and toward CSU & UC."),
       ("Faculty-driven, student-centered","Faculty own the academic recommendations; students are proactively identified and served.")]
x0,y0,w,h,gx,gy=0.6,2.42,3.92,1.83,0.185,0.2
dots=[NAVY,BLUE,GOLD,NAVY,BLUE,GOLD]
for i,(head,desc) in enumerate(princ):
    r,c=divmod(i,3); x=x0+c*(w+gx); y=y0+r*(h+gy)
    box(s3,x,y,w,h,fill=TINT,ln=LINE)
    box(s3,x+0.28,y+0.28,0.2,0.2,fill=dots[i],shape=MSO_SHAPE.OVAL)
    tf=textbox(s3,x+0.62,y+0.22,w-0.85,0.5)
    para(tf,[(head,14.5,NAVY,True)],first=True)
    tf=textbox(s3,x+0.28,y+0.72,w-0.55,1.0)
    para(tf,[(desc,11.5,DARK,False)],first=True,ls=1.04)

# reorder: move the newly-added slide (last) to right after slide 18 (index 17)
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
new_id = ids[-1]
after = ids[17]           # slide 18's sldId
sldIdLst.remove(new_id)
after.addnext(new_id)

# ---------- speaker notes ----------
NOTES = [
 ("[Slide — the charge]\nColleagues, here's what the one-time implementation funds are for. SB 135 and AB 135 establish Credit for Prior Learning as a systemwide initiative under the Master Plan for Career Education — and with that come new duties at every college: evaluate the prior learning and credentials of incoming students, adopt the credit recommendations our faculty discipline groups develop statewide, and accept CPL a student earned at another community college so the credit travels with them. The one-time funds help you stand that up locally — staffing, faculty engagement, student outreach, technology, employer partnerships, and room to innovate on new and stackable pathways. The key distinction for your planning: one-time dollars build local capacity; the ongoing dollars sustain the shared statewide backbone — so this is not a cliff.", s1),
 ("[Slide — the three priorities]\nThe Draft Implementation Funding model is outcomes-based — your college earns its share on measured CPL results, reflecting the statute's goal of advancing career attainment. Three priorities: Priority 1, Completion — growing certificate and degree completion through CPL awards. Priority 2, Access — reaching and enrolling more students and making CPL visible and accessible. Priority 3, Capacity and Mobility — building durable CPL that's documentable, interoperable, and portable across colleges and up to CSU and UC. I want to be clear this is a draft: the specific weighting and dollar allocations are still being worked out with the field. The framework — these three priorities — is what's settled.", s2),
 ("[Slide — guiding principles]\nFinally, the principles the model is built on — the fiscal design you'll care about. Outcomes-based, so funding follows results. Equitable and inclusive: a minimum-viable floor so even a small college can stand up a real program, a rural-college allowance, and support for our noncredit feeder institutions. Transparent and predictable — one published model where you can see how your allocation is derived. Sustainable — one-time to build, ongoing to sustain. Systemwide and portable — credit earned once travels. And faculty-driven and student-centered. Happy to take questions.", s3),
]
for txt, sld in NOTES:
    sld.notes_slide.notes_text_frame.text = txt

prs.save(OUT)
print("Saved:", OUT, "| total slides:", len(prs.slides._sldIdLst))
