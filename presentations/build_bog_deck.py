#!/usr/bin/env python3
"""Generate the CPL Initiative — Board of Governors update deck (Sam's 15-min segment).
Navy/white template matching Sam's sample slides. High-level, simple, board-appropriate.
Data as of live_metrics.json / fact_sheet_metrics.json 2026-07-16.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

REPO = "/home/user/cpl-project-tracker"
OUT = os.path.join(REPO, "presentations", "20260716_CPL_Initiative_BOG_Update.pptx")
SEAL = os.path.join(REPO, "cccco_seal.png")

# ---- palette (matched to Sam's samples) ----
NAVY   = RGBColor(0x15, 0x37, 0x63)   # deep navy — headers/title text
NAVY2  = RGBColor(0x0F, 0x2A, 0x4D)   # darker navy — cover bg
ACCENT = RGBColor(0x1D, 0x5C, 0x9E)   # medium blue — accent box/bar (the portal bar / blue box)
ACCENT2= RGBColor(0x2A, 0x6D, 0xB5)   # lighter blue
GOLD   = RGBColor(0xE8, 0xA3, 0x2C)   # gold accent (Get Started outline)
LIGHT  = RGBColor(0xEE, 0xF3, 0xF9)   # very light blue tile bg
LINE   = RGBColor(0xD5, 0xDF, 0xEA)   # tile border
GREY   = RGBColor(0x55, 0x63, 0x72)   # muted text
DARK   = RGBColor(0x22, 0x2A, 0x33)   # body text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x7D, 0x46)   # checkmarks

TITLE_FONT = "Cambria"   # serif, matches samples
BODY_FONT  = "Calibri"

EMW, EMH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMW
prs.slide_height = EMH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        # lightweight outer shadow
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (txt, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font, *rest) in para:
            italic = rest[0] if rest else False
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
            r.font.italic = italic
    return tb


def footer(s, note, pageno):
    text(s, 0.55, 7.06, 9.5, 0.35, [[(note, 9, GREY, False, BODY_FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 12.35, 7.06, 0.7, 0.35, [[(str(pageno), 10, GREY, False, BODY_FONT)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def stat_tile(s, x, y, w, h, number, label, num_color=NAVY, sub=None):
    rect(s, x, y, w, h, fill=LIGHT, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    paras = [[(number, 30, num_color, True, TITLE_FONT)]]
    if sub:
        paras.append([(sub, 10.5, ACCENT, True, BODY_FONT)])
    paras.append([(label, 11.5, GREY, False, BODY_FONT)])
    text(s, x + 0.15, y + 0.12, w - 0.3, h - 0.24, paras,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=0.98)


def title_block(s, kicker, title_lines, y=0.5):
    if kicker:
        text(s, 0.6, y, 12.1, 0.35, [[(kicker.upper(), 12.5, ACCENT, True, BODY_FONT)]])
        y += 0.42
    text(s, 0.6, y, 12.1, 1.0, [[(t, 30, NAVY, True, TITLE_FONT)] for t in title_lines],
         space_after=0, line_spacing=0.98)
    rect(s, 0.62, y + 0.02 + 0.62 * len(title_lines), 1.7, 0.06, fill=GOLD)
    return y + 0.62 * len(title_lines) + 0.28


# =====================================================================
# SLIDE 1 — COVER
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=NAVY2)
rect(s, 0, 6.95, 13.333, 0.55, fill=ACCENT)
if os.path.exists(SEAL):
    s.shapes.add_picture(SEAL, Inches(6.06), Inches(0.55), height=Inches(1.35))
text(s, 1.0, 2.15, 11.333, 0.5,
     [[("CREDIT FOR PRIOR LEARNING (CPL) INITIATIVE", 17, GOLD, True, BODY_FONT)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 2.7, 11.333, 1.5,
     [[("Update to the Board of Governors", 40, WHITE, True, TITLE_FONT)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 3.9, 11.333, 0.8,
     [[("Celebrating California’s 2026–27 Investment in Student Success", 19, RGBColor(0xCF,0xDD,0xEE), False, BODY_FONT, True)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 6.05, 11.333, 0.7,
     [[("Board of Governors Meeting  ·  July 2026", 14, WHITE, True, BODY_FONT)],
      [("California Community Colleges Chancellor’s Office", 12.5, RGBColor(0xB9,0xCB,0xE0), False, BODY_FONT)]],
     align=PP_ALIGN.CENTER, space_after=2)

# =====================================================================
# SLIDE 2 — POSITIONING STATEMENT (Sam's slide 1, recreated)
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
text(s, 0.85, 0.7, 11.633, 3.4,
     [[("California’s CPL Initiative is", 33, NAVY, True, TITLE_FONT)],
      [("Student-First, College-Supported, Employer-Engaged:", 33, NAVY, True, TITLE_FONT)],
      [("Moving From…", 33, NAVY, True, TITLE_FONT)]],
     align=PP_ALIGN.CENTER, space_after=4, line_spacing=1.0)
rect(s, 0.85, 4.15, 11.633, 2.75, fill=ACCENT)
text(s, 1.1, 4.35, 11.133, 2.35,
     [[("Passive  →  Proactive", 30, WHITE, False, TITLE_FONT)],
      [("Siloed  →  Shared", 30, WHITE, False, TITLE_FONT)],
      [("College  →  Career … & vice versa  :)", 30, WHITE, False, TITLE_FONT)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=6, line_spacing=1.0)

# =====================================================================
# SLIDE 3 — STATEWIDE IMPACT KPIs
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
title_block(s, "CPL in 2026", ["Statewide Momentum"])
tiles = [
    ("50,154", "Students served", NAVY, "since launch"),
    ("$326M", "Student tuition & time savings", GREEN, "to date"),
    ("12,679", "Credit recommendations", NAVY, "1,297 statewide · 11,382 local"),
    ("102", "Active colleges", ACCENT, "of 115 · 15 leading"),
    ("104K", "Units transcribed to records", NAVY, "236K identified as eligible"),
    ("$1.32B", "20-year economic impact", GREEN, "projected"),
]
x0, y0, tw, th, gx, gy = 0.6, 2.15, 3.98, 1.95, 0.19, 0.22
for i, (num, lab, col, sub) in enumerate(tiles):
    r, c = divmod(i, 3)
    stat_tile(s, x0 + c * (tw + gx), y0 + r * (th + gy), tw, th, num, lab, col, sub)
footer(s, "Source: MAP CPL Insights Dashboard · as of July 16, 2026", 3)

# =====================================================================
# SLIDE 4 — THE FUNDING WIN (theme centerpiece)
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=NAVY2)
rect(s, 0, 0, 13.333, 0.14, fill=GOLD)
text(s, 0.6, 0.55, 12.1, 0.4, [[("A VISION 2030 PRIORITY", 13, GOLD, True, BODY_FONT)]])
text(s, 0.6, 0.95, 12.1, 1.0, [[("2026–27 Budget: California Invests in CPL", 32, WHITE, True, TITLE_FONT)]])
# two big funding tiles
def fund_tile(x, big, unit, head, body):
    rect(s, x, 2.35, 5.85, 3.55, fill=RGBColor(0x18,0x3A,0x63), line=ACCENT2, line_w=1.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.35, 2.65, 5.15, 1.2,
         [[(big, 60, GOLD, True, TITLE_FONT), ("  " + unit, 22, RGBColor(0xCF,0xDD,0xEE), True, BODY_FONT)]])
    text(s, x+0.35, 3.95, 5.15, 0.5, [[(head, 18, WHITE, True, BODY_FONT)]])
    text(s, x+0.35, 4.5, 5.15, 1.25, [[(body, 14.5, RGBColor(0xCF,0xDD,0xEE), False, BODY_FONT)]],
         line_spacing=1.05)
fund_tile(0.6, "$7M", "ongoing", "Statewide infrastructure & operations",
          "Sustains the MAP platform, the statewide credit-recommendation library, and the staffing that keeps CPL running for every college.")
fund_tile(6.9, "$35M", "one-time", "Local CPL implementation",
          "Flows to colleges to build local CPL capacity — allocated by student outcomes, using the COBI Implementation Funding model.")
text(s, 0.6, 6.35, 12.1, 0.7,
     [[("Direct state support for a statewide system that turns work, service, and life experience into college credit — advancing Vision 2030 completion and equity goals.", 14, RGBColor(0xB9,0xCB,0xE0), False, BODY_FONT, True)]],
     align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 5 — THE 3 PRIORITIES OF THE COBI IMPLEMENTATION FUNDING MODEL
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
title_block(s, "COBI Implementation Funding model", ["Three Priorities Drive Every Dollar"])
text(s, 0.6, 2.02, 12.1, 0.5,
     [[("Each college earns its share of the $35M on measured outcomes across three priorities — modeled transparently in COBI.", 14.5, GREY, False, BODY_FONT, True)]])
# three big priority cards (official priority numbers, names, weights, descriptions)
pris = [
    ("PRIORITY 1", "30%", "Completion",
     "Increase CCC certificate & degree completion through CPL awards."),
    ("PRIORITY 2", "42%", "Access",
     "Increase college access through CPL — reaching more students."),
    ("PRIORITY 3", "28%", "Capacity & Mobility",
     "Increase CPL capacity, visibility, documentability, interoperability & mobility."),
]
x0, y0, tw, th, gx = 0.6, 2.62, 3.98, 2.85, 0.19
for i, (badge, pct, name, body) in enumerate(pris):
    x = x0 + i * (tw + gx)
    rect(s, x, y0, tw, th, fill=LIGHT, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # header band
    rect(s, x, y0, tw, 0.55, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y0+0.28, tw, 0.27, fill=ACCENT)  # square off header bottom
    text(s, x+0.25, y0, tw-0.4, 0.55, [[(badge, 14, WHITE, True, BODY_FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    # big weight
    text(s, x+0.25, y0+0.72, tw-0.5, 0.85,
         [[(pct, 44, NAVY, True, TITLE_FONT), ("  of the pool", 11.5, GREY, False, BODY_FONT)]])
    text(s, x+0.25, y0+1.62, tw-0.5, 0.4, [[(name, 18, ACCENT, True, BODY_FONT)]])
    text(s, x+0.25, y0+2.05, tw-0.5, 0.72, [[(body, 12.5, DARK, False, BODY_FONT)]], line_spacing=1.06)
# slim guardrails footer line
rect(s, 0.6, 5.72, 12.13, 0.72, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.9, 5.72, 11.55, 0.72,
     [[("Equity guardrails:  ", 13.5, GOLD, True, BODY_FONT),
       ("$150K minimum per college  ·  rural carve-out  ·  noncredit-feeder support  —  outcomes-based never means winner-take-all.", 13.5, WHITE, False, BODY_FONT)]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, "COBI → Implementation Funding tab · Chancellor-editable policy · every dollar computed live", 5)

# =====================================================================
# SLIDE 6 — STUDENT CPL PORTAL (Sam's slide 2, recreated)
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
text(s, 0.6, 0.45, 12.1, 0.4, [[("NEW · SOFT-LAUNCHED 2026", 13, GOLD, True, BODY_FONT)]])
text(s, 0.6, 0.82, 12.1, 0.6, [[("One Front Door for Every Californian", 30, NAVY, True, TITLE_FONT)]])
rect(s, 0.62, 1.5, 1.7, 0.06, fill=GOLD)
# portal mock
px, py, pw, ph = 0.6, 1.95, 12.13, 4.35
rect(s, px, py, pw, ph, fill=WHITE, line=LINE, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, px, py, pw, 0.7, fill=ACCENT)  # header bar
text(s, px+0.3, py, 6, 0.7, [[("CPL Student Portal", 16, WHITE, True, BODY_FONT)]], anchor=MSO_ANCHOR.MIDDLE)
rect(s, px+pw-2.05, py+0.16, 1.75, 0.4, fill=WHITE, line=GOLD, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, px+pw-2.05, py+0.16, 1.75, 0.4, [[("Get Started", 12, ACCENT, True, BODY_FONT)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, px+0.55, py+1.25, 7.2, 1.0,
     [[("Get Credit for Being ", 34, DARK, True, TITLE_FONT), ("You!", 34, ACCENT, True, TITLE_FONT)]])
text(s, px+0.55, py+2.25, 7.2, 1.2,
     [[("Credit for Prior Learning helps Californians turn what they’ve already learned — on the job, in the military, or through life — into college credit at California Community Colleges.", 15, GREY, False, BODY_FONT)]],
     line_spacing=1.12)
rect(s, px+0.55, py+3.5, 2.6, 0.55, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, px+0.55, py+3.5, 2.6, 0.55, [[("Create an Account  →", 14, WHITE, True, BODY_FONT)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# right panel accent
rect(s, px+8.4, py+1.1, 3.4, 2.9, fill=LIGHT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, px+8.6, py+1.1, 3.0, 2.9,
     [[("\U0001F393", 46, ACCENT, False, BODY_FONT)],
      [("Passive  →  Proactive", 16, NAVY, True, BODY_FONT)],
      [("Students find CPL — CPL no longer waits to be found.", 12.5, GREY, False, BODY_FONT, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=8, line_spacing=1.05)
text(s, 0.6, 6.5, 12.1, 0.5,
     [[("creditforbeingyou.org/main/student", 18, ACCENT, True, BODY_FONT)]],
     align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 7 — MY CPL STORY (Sam's slide 3, recreated as clean cards)
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
text(s, 0.6, 0.45, 7, 0.7, [[("My CPL Story", 32, NAVY, True, TITLE_FONT)]])
text(s, 7.2, 0.62, 5.55, 0.5, [[("map.rccd.edu/cplstories", 16, ACCENT, True, BODY_FONT)]],
     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 0.62, 1.2, 1.7, 0.06, fill=GOLD)
stories = [
    ("Luis C.", "26", "Marine Corps Veteran → College", "Military / JST", "San Bernardino Valley"),
    ("Nadine M.", "33.5", "Aviation Maintenance → A.S.", "Industry Certifications", "San Diego Miramar"),
    ("Eric O.", "27", "Auto Tech → B.S. (now faculty)", "Portfolio + Industry Certs", "Rio Hondo"),
    ("Sierra", "6", "Marine Veteran → Kinesiology", "Military / JST", "Cabrillo"),
    ("John N.", "29", "Firefighter → Public Safety B.S.", "Industry Certifications", "Moreno Valley"),
]
cw, gx = 2.3, 0.19
x0, y0, ch = 0.6, 1.55, 4.9
for i, (name, cr, path, typ, college) in enumerate(stories):
    x = x0 + i * (cw + gx)
    rect(s, x, y0, cw, ch, fill=WHITE, line=LINE, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # avatar circle with initials
    rect(s, x + cw/2 - 0.55, y0+0.3, 1.1, 1.1, fill=LIGHT, line=ACCENT, line_w=1.25, shape=MSO_SHAPE.OVAL)
    initials = "".join([p[0] for p in name.replace(".","").split()][:2])
    text(s, x + cw/2 - 0.55, y0+0.3, 1.1, 1.1, [[(initials, 26, ACCENT, True, TITLE_FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.1, y0+1.5, cw-0.2, 0.4, [[(name, 15, NAVY, True, BODY_FONT)]], align=PP_ALIGN.CENTER)
    # credits badge
    rect(s, x+cw/2-0.95, y0+1.92, 1.9, 0.42, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+cw/2-0.95, y0+1.92, 1.9, 0.42, [[(cr + " credits", 12, WHITE, True, BODY_FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.12, y0+2.5, cw-0.24, 1.4,
         [[("✓ " + typ, 11, GREEN, True, BODY_FONT)],
          [(path, 12, DARK, False, BODY_FONT)]],
         align=PP_ALIGN.CENTER, space_after=6, line_spacing=1.02)
    rect(s, x+0.12, y0+ch-0.55, cw-0.24, 0.44, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.12, y0+ch-0.55, cw-0.24, 0.44, [[(college, 10.5, GREY, True, BODY_FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Real students · real credit · real momentum toward a degree", 7)

# =====================================================================
# SLIDE 8 — VIDEO: EMT at Moreno Valley
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=NAVY2)
text(s, 0.6, 0.55, 12.1, 0.4, [[("STUDENT STORY · 2-MINUTE VIDEO", 13, GOLD, True, BODY_FONT)]])
text(s, 0.6, 0.98, 12.1, 0.7, [[("From First Responder to a Degree", 30, WHITE, True, TITLE_FONT)]])
# video frame
vx, vy, vw, vh = 2.15, 2.0, 9.0, 4.05
rect(s, vx, vy, vw, vh, fill=RGBColor(0x0A,0x1F,0x38), line=ACCENT2, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, vx+vw/2-0.7, vy+vh/2-0.7, 1.4, 1.4, fill=GOLD, shape=MSO_SHAPE.OVAL)
tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(vx+vw/2-0.24), Inches(vy+vh/2-0.32), Inches(0.55), Inches(0.64))
tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = NAVY2; tri.line.fill.background(); tri.shadow.inherit=False
text(s, vx, vy+vh-1.05, vw, 0.9,
     [[("An EMT’s path — Moreno Valley College", 17, WHITE, True, BODY_FONT)],
      [("[ Insert 2-minute student story video here ]", 12.5, RGBColor(0xB9,0xCB,0xE0), False, BODY_FONT, True)]],
     align=PP_ALIGN.CENTER, space_after=3)
text(s, 0.6, 6.35, 12.1, 0.5,
     [[("Emergency medical training → college credit → a faster path to a bachelor’s degree.", 14.5, RGBColor(0xCF,0xDD,0xEE), False, BODY_FONT, True)]],
     align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 9 — CPL IN ACTION: Moreno Valley Emergency Management B.S. (illustrative pathway)
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
title_block(s, "CPL in action · a sample pathway", ["Moreno Valley College — B.S. Emergency Management"])
text(s, 0.6, 2.0, 12.1, 0.5,
     [[("Illustrative — how a first responder stacks prior learning into a bachelor’s degree.", 14.5, GREY, False, BODY_FONT, True)]])
# 4-step flow
steps = [
    ("PRIOR LEARNING", "EMT certification · Fire academy · EMS & first-responder training · military service", ACCENT),
    ("CPL CREDIT", "Verified against statewide credit recommendations — e.g. EMT (7u) + Basic Fire Academy (up to 24u)", ACCENT2),
    ("QUALIFYING A.S.", "CPL applies to the associate degree that gates entry to the B.S.", NAVY),
    ("B.S. EMERGENCY MGMT", "Moreno Valley College · approved March 2026 — reached faster & cheaper", GREEN),
]
bx, by, bw, bh, gap = 0.6, 2.75, 2.86, 2.5, 0.24
for i, (head, body, col) in enumerate(steps):
    x = bx + i * (bw + gap)
    rect(s, x, by, bw, bh, fill=LIGHT, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, by, bw, 0.5, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, by+0.25, bw, 0.25, fill=col)  # square off bottom of header
    text(s, x+0.1, by, bw-0.2, 0.5, [[(head, 12.5, WHITE, True, BODY_FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.18, by+0.62, bw-0.36, bh-0.75, [[(body, 12.5, DARK, False, BODY_FONT)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.06)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+bw+0.01), Inches(by+bh/2-0.16), Inches(gap-0.02), Inches(0.32))
        ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background(); ar.shadow.inherit=False
rect(s, 0.6, 5.55, 12.13, 0.95, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.9, 5.68, 11.5, 0.7,
     [[("Because CPL counts toward the qualifying associate degree, students enter the bachelor’s program with credit already in hand — saving time, tuition, and repeated coursework.", 14, WHITE, False, BODY_FONT)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s, "Illustrative pathway · unit values reflect published statewide Fire/EMS credit recommendations", 9)

# =====================================================================
# SLIDE 10 — PARTNERSHIPS
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
title_block(s, "College-Supported, Employer-Engaged", ["Built With — Not For — the Field"])
parts = [
    ("Faculty & ASCCC", "Discipline faculty workgroups author the statewide credit recommendations colleges adopt."),
    ("Employers & Industry", "Industry certifications (ASE, Cal-JAC & more) map directly to college credit."),
    ("Cal OES / CSTI", "New collaboration to articulate CSTI certifications into CCC Fire & Public Safety courses."),
    ("The Military", "Joint Services Transcript (JST) credit for service members and veterans."),
    ("Apprenticeship", "Registered apprenticeship hours recognized toward degrees & certificates."),
    ("CAEL / ACE", "National standards & assessment partners anchoring credit quality."),
]
x0, y0, tw, th, gx, gy = 0.6, 2.2, 3.98, 1.95, 0.19, 0.2
for i, (head, body) in enumerate(parts):
    r, c = divmod(i, 3)
    x = x0 + c*(tw+gx); y = y0 + r*(th+gy)
    rect(s, x, y, tw, th, fill=WHITE, line=LINE, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, tw, 0.1, fill=ACCENT)
    text(s, x+0.24, y+0.25, tw-0.45, th-0.4,
         [[(head, 15.5, NAVY, True, BODY_FONT)],
          [(body, 12.5, DARK, False, BODY_FONT)]],
         space_after=5, line_spacing=1.06)
footer(s, "Siloed → Shared: credit recommendations built once, adopted statewide", 10)

# =====================================================================
# SLIDE 11 — HAND-OFF: FIRE TECHNOLOGY (tee up Miramar)
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=NAVY2)
rect(s, 0, 0, 13.333, 0.14, fill=GOLD)
text(s, 0.6, 0.55, 12.1, 0.4, [[("A STATEWIDE MODEL IN ACTION", 13, GOLD, True, BODY_FONT)]])
text(s, 0.6, 0.98, 12.1, 0.7, [[("Fire Technology — Leading the Way", 32, WHITE, True, TITLE_FONT)]])
fstats = [("294", "statewide Fire credit recommendations"),
          ("342", "college adoptions to date"),
          ("45", "colleges eligible to adopt"),
          ("~40", "colleges with SFT programs & academies")]
x0, y0, tw, th, gx = 0.6, 2.05, 2.95, 1.65, 0.19
for i,(num,lab) in enumerate(fstats):
    x = x0 + i*(tw+gx)
    rect(s, x, y0, tw, th, fill=RGBColor(0x18,0x3A,0x63), line=ACCENT2, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.12, y0+0.15, tw-0.24, th-0.3,
         [[(num, 34, GOLD, True, TITLE_FONT)], [(lab, 12, WHITE, False, BODY_FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=3, line_spacing=0.98)
# new collaboration callout
rect(s, 0.6, 4.0, 12.13, 1.4, fill=RGBColor(0x18,0x3A,0x63), line=GOLD, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.9, 4.14, 11.5, 0.4, [[("NEW — Cal OES / CSTI collaboration", 15, GOLD, True, BODY_FONT)]])
text(s, 0.9, 4.55, 11.55, 0.8,
     [[("Articulating CSTI certifications and offering stackable CSTI credit within existing Fire & Public Safety courses — expanding sector capacity and helping California “train the trainers.”", 14.5, WHITE, False, BODY_FONT)]],
     line_spacing=1.08)
# handoff to speakers
rect(s, 0.6, 5.6, 12.13, 1.15, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.9, 5.72, 11.5, 0.9,
     [[("Next: San Diego Miramar College", 18, WHITE, True, BODY_FONT)],
      [("Faculty & Dean — leading the statewide Fire Technology workgroup, in their own words.", 14.5, RGBColor(0xE7,0xEF,0xF8), False, BODY_FONT, True)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=3)

# =====================================================================
# SLIDE 12 — CLOSE / THANK YOU / Q&A
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=NAVY2)
rect(s, 0, 6.95, 13.333, 0.55, fill=ACCENT)
if os.path.exists(SEAL):
    s.shapes.add_picture(SEAL, Inches(6.16), Inches(0.7), height=Inches(1.15))
text(s, 1.0, 2.25, 11.333, 0.9,
     [[("Student-First. College-Supported. Employer-Engaged.", 26, WHITE, True, TITLE_FONT)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 3.25, 11.333, 0.6,
     [[("Thank You", 40, GOLD, True, TITLE_FONT)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 4.25, 11.333, 0.6,
     [[("Questions & Discussion", 20, RGBColor(0xCF,0xDD,0xEE), False, BODY_FONT, True)]],
     align=PP_ALIGN.CENTER)
text(s, 1.0, 6.05, 11.333, 0.7,
     [[("map.rccd.edu  ·  creditforbeingyou.org  ·  MAP@rccd.edu", 14, WHITE, True, BODY_FONT)]],
     align=PP_ALIGN.CENTER)

# =====================================================================
# SPEAKER NOTES — a talking script in each slide's notes field (~13 min
# of Sam speaking + the 2-min video = the first 15 minutes).
# =====================================================================
NOTES = [
    # 1 — Cover (~20s)
    "[~20 sec — open]\n"
    "Good morning, and thank you. I'm excited to give the Board a quick update on the Credit for Prior Learning Initiative — and to celebrate a real milestone. This year, the state made a direct investment in CPL as part of California's Vision 2030 commitment. I'll take about 15 minutes on where the initiative stands, then hand off to our Fire Technology faculty leaders from San Diego Miramar, and we'll leave 15 minutes for your questions.",
    # 2 — Positioning (~60s)
    "[~60 sec — the vision in one slide]\n"
    "Here's the whole philosophy on one slide. California's CPL Initiative is student-first, college-supported, and employer-engaged. And it's about a shift in posture.\n"
    "• Passive to proactive — we no longer wait for a student to know to ask for credit; the system reaches out to them.\n"
    "• Siloed to shared — a faculty workgroup writes a credit recommendation once, and colleges across the state adopt it, instead of 116 colleges solving the same problem alone.\n"
    "• And college to career — and, just as importantly, career back to college. We honor what people already learned on the job, in the military, and in life, and we turn it into a credential.\n"
    "Everything I show you next is an example of one of those three shifts.",
    # 3 — KPIs (~90s)
    "[~90 sec — the numbers]\n"
    "First, the scale. To date, CPL has served more than 50,000 students across California — roughly half military and veterans, half workforce and other learners. Those students have identified 236,000 units of eligible prior learning, and we've transcribed 104,000 of those units onto student records — real credit, on real transcripts.\n"
    "That translates to about $326 million in tuition and time savings for students so far — and a projected $1.3 billion in economic impact over 20 years.\n"
    "Underneath that: nearly 12,700 credit recommendations, and 102 active colleges out of 115, with 15 leading colleges already at scale. This is no longer a pilot — it's statewide infrastructure. All of these numbers are live on our dashboard.",
    # 4 — Funding win (~90s)
    "[~90 sec — THE celebration]\n"
    "And this is what we're here to celebrate. In the 2026–27 budget, California invested directly in CPL as a Vision 2030 priority — two pieces.\n"
    "First, $7 million ongoing to sustain the statewide infrastructure and operations — the MAP platform, the shared credit-recommendation library, and the staffing that keeps this running for every college. That's the part that makes CPL a permanent system, not a grant that sunsets.\n"
    "Second, $35 million one-time to fund local implementation — money that flows to colleges to build their own CPL capacity. And critically, that money is allocated by student outcomes, using the COBI Implementation Funding model. Let me show you how that works.",
    # 5 — COBI 3 priorities (~80s)
    "[~80 sec — the three priorities]\n"
    "The $35 million doesn't go out as a flat check. Every college earns its share based on measured outcomes across three priorities — and it's fully transparent in our COBI Implementation Funding model, where you can see every dollar computed live.\n"
    "Priority 1 is completion — 30% of the pool rewards actual certificates and degrees earned through CPL. That's the ultimate goal: not just credit, but credentials.\n"
    "Priority 2 is access — the largest share at 42% — because the first job is simply reaching more students and getting their prior learning recognized.\n"
    "And Priority 3, 28%, is capacity and mobility — building the systems that make that credit visible, documented, and portable from college to college.\n"
    "And we wrapped it in equity guardrails: a $150,000 minimum so small colleges can stand up a real program, a rural carve-out, and support for our noncredit feeder institutions. Outcomes-based — but no college left behind.",
    # 6 — Portal (~60s)
    "[~60 sec — passive to proactive, made real]\n"
    "Here's 'passive to proactive' made real. This year we soft-launched the new CPL Student Portal — creditforbeingyou.org. It's one front door for every Californian: 'Get Credit for Being You.'\n"
    "A student — a veteran, a working parent, a career-changer — creates an account, tells us what they've done, and the system helps them turn that experience into college credit. Before, credit for prior learning was something you had to know existed and know to ask for. Now it comes to the student. This is early, but it's the model for scale.",
    # 7 — Stories (~55s)
    "[~55 sec — the human proof]\n"
    "Behind every number is a person. These are real students from our My CPL Story collection. Luis — a Marine sergeant — got 26 credits from his Joint Services Transcript. Nadine saved a full year toward her aviation maintenance degree with her industry certifications. Eric, an auto technician, earned 27 credits and is now teaching at Rio Hondo. Sierra is applying her military first-aid training to a kinesiology degree. And John — a firefighter — earned 29 credits toward a public safety degree.\n"
    "Different fields, different backgrounds, same story: what they already knew counted.",
    # 8 — Video (~2 min video)
    "[~2 min — PLAY VIDEO]\n"
    "I want you to hear this directly from a student. This is a two-minute story of an EMT at Moreno Valley College — someone whose emergency medical training became college credit and a faster path to a degree.\n"
    "[Click play. Let the video run. No need to narrate.]\n"
    "ACTION ITEM: drop the 2-minute EMT video onto this slide (Insert → Video) before presenting.",
    # 9 — Pathway (~75s)
    "[~75 sec — how it stacks]\n"
    "Let me show you how that EMT's credit actually stacks into a degree. This is illustrative, using Moreno Valley's new Bachelor of Science in Emergency Management, approved this past March.\n"
    "It starts with prior learning — EMT certification, a fire academy, EMS and first-responder training, military service. That maps against our statewide credit recommendations into real CPL credit — for example, 7 units for EMT, up to 24 for a basic fire academy. That credit applies to the qualifying associate degree that gates entry to the bachelor's. So the student walks into the B.S. program with credit already in hand — reaching the degree faster and cheaper, without repeating what they already know. That's CPL doing exactly what it's designed to do.",
    # 10 — Partnerships (~60s)
    "[~60 sec — built WITH the field]\n"
    "None of this works if it's built in a vacuum. CPL is built with the field. Our discipline faculty and the Academic Senate author the credit recommendations — faculty own the academic judgment. Employers and industry bring the certifications — ASE, the apprenticeship trusts, and more — that map to credit. The military brings the Joint Services Transcript. And national partners like CAEL and ACE anchor the quality.\n"
    "The newest one is on the next slide, and it's a big one — a collaboration with Cal OES and the state fire training system. Which is the perfect place to hand off.",
    # 11 — Fire handoff (~75s)
    "[~75 sec — set up Miramar, then hand off]\n"
    "Fire Technology is our flagship example of the whole model working. A statewide faculty workgroup produced 294 Fire credit recommendations, which have already been adopted 342 times, with 45 colleges eligible to adopt — out of roughly 40 colleges statewide running state-fire-training programs and academies.\n"
    "And the new frontier: a collaboration with Cal OES and CSTI to articulate CSTI certifications and offer stackable CSTI credit inside colleges' existing Fire and Public Safety courses. That expands the state's capacity to train first responders — and to train the trainers.\n"
    "The people who led that workgroup are here. So let me hand it to our faculty and dean from San Diego Miramar College, who will speak from their own experience.\n"
    "[Hand off. Their segment is ~15 min, no slides. This slide can stay up.]",
    # 12 — Close (~20s)
    "[~20 sec — if you close the segment]\n"
    "Thank you again. Student-first, college-supported, employer-engaged — that's California's CPL Initiative, and with this year's investment it's built to last. We're happy to take your questions.",
]

for i, sld in enumerate(prs.slides):
    if i < len(NOTES):
        sld.notes_slide.notes_text_frame.text = NOTES[i]

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
print("Notes attached:", len(NOTES))
