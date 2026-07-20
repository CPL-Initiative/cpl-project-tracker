#!/usr/bin/env python3
"""CBO overview — 3 slides on how the $35M one-time CPL funds support local
implementation: the Draft Implementation Funding model's 3 priorities + guiding
principles. High-level, NO model dollar amounts (still in discussion). Same
navy/white template as the BOG deck. Speaker notes on each slide."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = "/home/user/cpl-project-tracker"
OUT = os.path.join(REPO, "presentations", "20260720_CPL_CBO_Implementation_Funding.pptx")
SEAL = os.path.join(REPO, "cccco_seal.png")

NAVY   = RGBColor(0x15, 0x37, 0x63)
NAVY2  = RGBColor(0x0F, 0x2A, 0x4D)
ACCENT = RGBColor(0x1D, 0x5C, 0x9E)
ACCENT2= RGBColor(0x2A, 0x6D, 0xB5)
GOLD   = RGBColor(0xE8, 0xA3, 0x2C)
LIGHT  = RGBColor(0xEE, 0xF3, 0xF9)
LINE   = RGBColor(0xD5, 0xDF, 0xEA)
GREY   = RGBColor(0x55, 0x63, 0x72)
DARK   = RGBColor(0x22, 0x2A, 0x33)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x7D, 0x46)
TITLE_FONT = "Cambria"
BODY_FONT  = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0); p.line_spacing = line_spacing
        for (t, sz, col, b, fn, *rest) in para:
            it = rest[0] if rest else False
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = col; r.font.name = fn; r.font.italic = it
    return tb

def title_block(s, kicker, title_lines, y=0.5):
    if kicker:
        text(s, 0.6, y, 12.1, 0.35, [[(kicker.upper(), 12.5, ACCENT, True, BODY_FONT)]]); y += 0.42
    text(s, 0.6, y, 12.1, 1.0, [[(t, 29, NAVY, True, TITLE_FONT)] for t in title_lines], space_after=0, line_spacing=0.98)
    rect(s, 0.62, y + 0.02 + 0.6*len(title_lines), 1.7, 0.06, fill=GOLD)
    return y + 0.6*len(title_lines) + 0.26

def footer(s, note, pageno):
    text(s, 0.6, 7.06, 9.8, 0.35, [[(note, 9.5, GREY, False, BODY_FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 12.3, 7.06, 0.7, 0.35, [[(str(pageno), 10, GREY, False, BODY_FONT)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def checklist(s, x, y, w, head, items, head_color=ACCENT):
    text(s, x, y, w, 0.4, [[(head, 15.5, NAVY, True, BODY_FONT)]])
    yy = y + 0.5
    for it in items:
        rect(s, x+0.02, yy+0.06, 0.16, 0.16, fill=head_color, shape=MSO_SHAPE.OVAL)
        text(s, x+0.32, yy-0.02, w-0.34, 0.6, [[(it, 13.5, DARK, False, BODY_FONT)]], line_spacing=1.03)
        # measure rough height by line count
        lines = max(1, int(len(it)/46)+1)
        yy += 0.34 + 0.24*(lines-1)
    return yy

# =====================================================================
# SLIDE 1 — What the one-time funds do (local implementation)
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
title_block(s, "$35M one-time · local college implementation", ["Standing Up CPL at Every College"])
text(s, 0.6, 1.95, 12.1, 0.5,
     [[("SB 135 / AB 135 establishes CPL as a ", 14.5, GREY, False, BODY_FONT, True),
       ("systemwide initiative", 14.5, NAVY, True, BODY_FONT),
       (" under the Master Plan for Career Education. One-time funds help each college build local capacity to meet new duties.", 14.5, GREY, False, BODY_FONT, True)]])
# two columns
lx, rx, cw, cy = 0.6, 6.95, 5.78, 2.7
rect(s, lx, cy, cw, 3.0, fill=LIGHT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, rx, cy, cw, 3.0, fill=LIGHT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, lx, cy, cw, 0.1, fill=ACCENT); rect(s, rx, cy, cw, 0.1, fill=GOLD)
checklist(s, lx+0.35, cy+0.32, cw-0.6, "What colleges now do", [
    "Evaluate incoming students’ prior learning & credentials for credit",
    "Adopt faculty-developed statewide credit recommendations",
    "Accept transcribed CPL from other colleges — credit that travels",
    "Build local capacity: people, processes & technology",
])
checklist(s, rx+0.35, cy+0.32, cw-0.6, "What the funds seed", [
    "Dedicated CPL staffing & faculty engagement",
    "Student identification, outreach & advising",
    "Employer & industry partnerships",
    "New & stackable pathways — local CPL innovation",
], head_color=GOLD)
rect(s, 0.6, 5.95, 12.13, 0.95, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.9, 6.06, 11.55, 0.75,
     [[("One-time funds build ", 14, WHITE, False, BODY_FONT),
       ("local capacity", 14, GOLD, True, BODY_FONT),
       ("; ongoing funds sustain the ", 14, WHITE, False, BODY_FONT),
       ("shared statewide backbone", 14, GOLD, True, BODY_FONT),
       (" — the technology and faculty credit recommendations every college draws on.", 14, WHITE, False, BODY_FONT)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
footer(s, "Source: SB 111 (appropriation) · SB 135 / AB 135 (Higher Education budget trailer bill) · Master Plan for Career Education", 1)

# =====================================================================
# SLIDE 2 — Draft Implementation Funding model: 3 priorities
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
title_block(s, "Draft Implementation Funding model", ["Funds Follow Outcomes — Across Three Priorities"])
text(s, 0.6, 2.0, 12.1, 0.5,
     [[("Each college’s share is earned on measured CPL outcomes — reflecting the statute’s goal of advancing career attainment through CPL.", 14.5, GREY, False, BODY_FONT, True)]])
pris = [
    ("PRIORITY 1", "Completion", "Grow certificate & degree completion — and career attainment — through CPL awards."),
    ("PRIORITY 2", "Access", "Reach and enroll more students through CPL, and make it visible and accessible."),
    ("PRIORITY 3", "Capacity & Mobility", "Build durable CPL — documentable, interoperable, and portable across colleges and segments."),
]
x0, y0, tw, th, gx = 0.6, 2.65, 3.98, 2.75, 0.19
for i, (badge, name, body) in enumerate(pris):
    x = x0 + i*(tw+gx)
    rect(s, x, y0, tw, th, fill=LIGHT, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y0, tw, 0.55, fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y0+0.28, tw, 0.27, fill=ACCENT)
    text(s, x+0.25, y0, tw-0.4, 0.55, [[(badge, 14, WHITE, True, BODY_FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.28, y0+0.78, tw-0.5, 0.7, [[(name, 24, NAVY, True, TITLE_FONT)]])
    text(s, x+0.28, y0+1.62, tw-0.55, 1.0, [[(body, 13.5, DARK, False, BODY_FONT)]], line_spacing=1.08)
rect(s, 0.6, 5.68, 12.13, 0.72, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.9, 5.68, 11.55, 0.72,
     [[("DRAFT: ", 13.5, GOLD, True, BODY_FONT),
       ("priority weighting and dollar allocations are still in development with the field — the framework is what’s settled.", 13.5, WHITE, False, BODY_FONT)]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, "COBI Implementation Funding model · outcomes measured per college via the MAP platform", 2)

# =====================================================================
# SLIDE 3 — Guiding principles
# =====================================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, fill=WHITE)
title_block(s, "How the model is designed", ["Guiding Principles"])
princ = [
    ("Outcomes-based", "Funding follows measured results — not headcount alone."),
    ("Equitable & inclusive", "A minimum-viable floor so every participating college can build a real program; a rural-college allowance; noncredit-feeder support."),
    ("Transparent & predictable", "One clear, published model — colleges can see how their allocation is derived."),
    ("Sustainable", "One-time funds build local capacity; ongoing funds sustain the statewide backbone."),
    ("Systemwide & portable", "Credit earned once travels — aligned across colleges and toward CSU & UC."),
    ("Faculty-driven, student-centered", "Faculty own the academic recommendations; students are proactively identified and served."),
]
x0, y0, tw, th, gx, gy = 0.6, 2.15, 3.98, 1.95, 0.19, 0.2
for i, (head, body) in enumerate(princ):
    r, c = divmod(i, 3)
    x = x0 + c*(tw+gx); y = y0 + r*(th+gy)
    rect(s, x, y, tw, th, fill=WHITE, line=LINE, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, 0.12, th, fill=GOLD if r==1 else ACCENT)
    text(s, x+0.32, y+0.22, tw-0.5, th-0.4,
         [[(head, 15.5, NAVY, True, BODY_FONT)],
          [(body, 12.5, DARK, False, BODY_FONT)]],
         space_after=5, line_spacing=1.05)
footer(s, "Draft Implementation Funding model — guiding principles · California Community Colleges · CPL Initiative", 3)

# =====================================================================
# SPEAKER NOTES
# =====================================================================
NOTES = [
    "[Slide 1 — the charge]\n"
    "Colleagues, here's what the one-time implementation funds are for. SB 135 and AB 135 establish Credit for Prior Learning as a systemwide initiative under the Master Plan for Career Education — and with that come new duties at every college: evaluate the prior learning and credentials of incoming students, adopt the credit recommendations our faculty discipline groups develop statewide, and accept CPL that students earned at another community college so the credit travels with them.\n"
    "The one-time funds help you stand that up locally — staffing, faculty engagement, student outreach, technology, employer partnerships, and room to innovate on new and stackable pathways. The key distinction for your planning: one-time dollars build local capacity; the ongoing dollars sustain the shared statewide backbone — the technology and the faculty credit recommendations you all draw on — so this isn't a cliff.",

    "[Slide 2 — the three priorities]\n"
    "The Draft Implementation Funding model is outcomes-based — your college earns its share on measured CPL results, reflecting the statute's goal of advancing career attainment. There are three priorities.\n"
    "Priority 1, Completion — growing certificate and degree completion, and career attainment, through CPL awards. Priority 2, Access — reaching and enrolling more students through CPL and making it visible and accessible. Priority 3, Capacity and Mobility — building durable CPL that's documentable, interoperable, and portable across colleges and up to CSU and UC.\n"
    "I want to be clear this is a draft: the specific weighting and dollar allocations are still being worked out with the field. The framework — these three priorities — is what's settled, and that's what I'm sharing today.",

    "[Slide 3 — guiding principles]\n"
    "Finally, the principles the model is built on — this is really the fiscal design you'll care about. It's outcomes-based, so funding follows results. It's equitable and inclusive: a minimum-viable floor so even a small college can stand up a real program, a rural-college allowance, and support for our noncredit feeder institutions. It's transparent and predictable — one published model where you can see how your allocation is derived. It's sustainable — one-time to build, ongoing to sustain. It's systemwide and portable — credit earned once travels. And it's faculty-driven and student-centered — faculty own the academic judgment, and students are proactively identified and served.\n"
    "Happy to take questions on any of this.",
]
for i, sld in enumerate(prs.slides):
    if i < len(NOTES): sld.notes_slide.notes_text_frame.text = NOTES[i]

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("Saved:", OUT, "| slides", len(prs.slides._sldIdLst))
