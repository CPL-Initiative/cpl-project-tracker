from pptx import Presentation
from pptx.util import Pt
import notes as N

import sys
SRC = sys.argv[1] if len(sys.argv) > 1 else '20260810_CAC_CPL_Presentation_Pathways_Animated.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else '20260810_CAC_CPL_Presentation_Pathways_Animated_Notes.pptx'

prs = Presentation(SRC)
applied = 0
for i, sl in enumerate(prs.slides):
    if i not in N.NOTES:
        continue
    body = N.NOTES[i].strip('\n')
    if i == 5:                                    # spine slide carries the background
        body += '\n' + N.PRE_APPRENTICESHIP.rstrip()
    existing = ''
    if sl.has_notes_slide:
        existing = (sl.notes_slide.notes_text_frame.text or '').strip()
    # slide 10 (index 9) carries "Video - Need Audio" -- already folded into the script
    if existing and existing.lower() not in body.lower() and 'need audio' not in existing.lower():
        body += '\n\n[ORIGINAL NOTE FROM THE DECK: %s]' % existing
    tf = sl.notes_slide.notes_text_frame
    tf.text = body
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(11)
            r.font.name = 'Calibri'
    applied += 1
prs.save(OUT)
print('wrote', OUT, '| notes on', applied, 'slides')
