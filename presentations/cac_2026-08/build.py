# -*- coding: utf-8 -*-
"""Rebuild slides 6/7/8 of the CAC deck as calm, click-built pathway slides.

Adds:  spine slide + 3 trade slides (replacing 6/7/8) + 2 appendix crosswalk slides.
Keeps: the deck's own master, theme, CCC logo and slide numbering.
"""
import sys, uuid
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
import spec as S

NS = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
      'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# ------------------------------------------------------------------ geometry
SW, SH = 13.333, 7.5
ML = 0.40
CW = SW - 2 * ML                          # 12.533

COLX = [0.40, 2.875, 5.35, 8.355, 10.83]
COLW = [2.10, 2.10, 2.63, 2.10, 2.10]

KICK_Y, TITLE_Y, SUB_Y = 0.30, 0.52, 1.06
HDR_Y, HDR_H = 1.46, 0.62
BOD_Y, BOD_H = 2.14, 3.02                 # -> 5.16
CHIP_Y, CHIP_H = 5.40, 0.46               # -> 5.86  (master rule sits at 6.23)

LANE_H   = 1.32
LANE_A_Y = 2.14                           # -> 3.46
LANE_B_Y = 3.84                           # -> 5.16


# ------------------------------------------------------------------ helpers
def tint(hexstr, amount):
    r, g, b = (int(hexstr[i:i + 2], 16) for i in (0, 2, 4))
    f = lambda c: min(255, max(0, int(round(c + (255 - c) * amount))))
    return '%02X%02X%02X' % (f(r), f(g), f(b))


def shade(hexstr, amount):
    r, g, b = (int(hexstr[i:i + 2], 16) for i in (0, 2, 4))
    f = lambda c: min(255, max(0, int(round(c * (1 - amount)))))
    return '%02X%02X%02X' % (f(r), f(g), f(b))


def _plain(shp):
    shp.shadow.inherit = False


def rrect(sl, x, y, w, h, fill=None, line=None, radius=0.07, lw=1.0):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    half = min(w, h) / 2.0
    s.adjustments[0] = max(0.0, min(0.5, radius / half)) if half else 0.0
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = RGBColor.from_string(line); s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    _plain(s)
    return s


def rect(sl, x, y, w, h, fill):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    s.line.fill.background()
    _plain(s)
    return s


def oval(sl, x, y, d, fill):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    s.line.fill.background()
    _plain(s)
    return s


def chevron(sl, x, y, w, h, fill):
    s = sl.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    s.line.fill.background()
    _plain(s)
    return s


def star(sl, x, y, d, fill):
    s = sl.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(fill)
    s.line.fill.background()
    _plain(s)
    return s


def _spc(run, val):
    run.font._rPr.set('spc', str(int(val)))


def tbox(sl, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, d in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = d.get('align', align)
        if d.get('before') is not None:
            p.space_before = Pt(d['before'])
        if d.get('after') is not None:
            p.space_after = Pt(d['after'])
        if d.get('ls'):
            p.line_spacing = d['ls']
        for rd in (d.get('runs') or [d]):
            r = p.add_run()
            r.text = rd.get('text', '')
            f = r.font
            f.name = rd.get('font', S.BODY_FONT)
            f.size = Pt(rd.get('size', 10))
            f.bold = bool(rd.get('bold'))
            f.italic = bool(rd.get('italic'))
            f.color.rgb = RGBColor.from_string(rd.get('color', S.GREY))
            if rd.get('spc'):
                _spc(r, rd['spc'])
    return tb


# ------------------------------------------------------------------ scaffold
def add_slide_number(sl):
    tb = tbox(sl, 12.10, 6.70, 0.78, 0.30,
              [{'text': '#', 'size': 10, 'color': S.GREY_LT, 'align': PP_ALIGN.RIGHT}])
    p = tb.text_frame.paragraphs[0]
    r = p.runs[0]
    fld = etree.SubElement(p._p, qn('a:fld'))
    fld.set('id', '{%s}' % str(uuid.uuid4()).upper())
    fld.set('type', 'slidenum')
    rPr = r._r.find(qn('a:rPr'))
    if rPr is not None:
        fld.append(deepcopy(rPr))
    t = etree.SubElement(fld, qn('a:t')); t.text = '2'
    r._r.addprevious(fld)
    r._r.getparent().remove(r._r)
    return tb


def new_slide(prs, layout):
    sl = prs.slides.add_slide(layout)
    for ph in list(sl.placeholders):
        if ph.placeholder_format.type != PP_PLACEHOLDER.SLIDE_NUMBER:
            ph._element.getparent().remove(ph._element)
    add_slide_number(sl)
    return sl


def add_heading(sl, kicker, title, sub, title_size=27):
    tbox(sl, ML, KICK_Y, CW, 0.22,
         [{'text': kicker.upper(), 'size': 9.5, 'bold': True, 'color': S.CCC_BLUE, 'spc': 90}])
    tbox(sl, ML, TITLE_Y, CW, 0.52,
         [{'text': title, 'size': title_size, 'bold': True, 'color': S.NAVY, 'font': S.TITLE_FONT}])
    tbox(sl, ML, SUB_Y, CW - 0.2, 0.34,
         [{'text': sub, 'size': 11, 'color': S.GREY, 'ls': 1.0}])


def stage_header(sl, i, x, w, label=None):
    c = S.STAGE_COLORS[i]
    bar = rrect(sl, x, HDR_Y, w, HDR_H, fill=c, radius=0.09)
    d = 0.30
    dy = HDR_Y + (HDR_H - d) / 2
    disc = oval(sl, x + 0.14, dy, d, S.WHITE)
    num = tbox(sl, x + 0.14, dy + 0.03, d, d - 0.04,
               [{'text': str(i + 1), 'size': 11, 'bold': True, 'color': c,
                 'align': PP_ALIGN.CENTER}])
    lab = tbox(sl, x + 0.50, HDR_Y + 0.04, w - 0.60, HDR_H - 0.08,
               [{'text': (label or S.STAGE_NAMES[i]), 'size': 9, 'bold': True,
                 'color': S.WHITE, 'ls': 0.9, 'spc': 20}], anchor=MSO_ANCHOR.MIDDLE)
    return [bar, disc, num, lab]


def body_card(sl, i, x, w, y=BOD_Y, h=BOD_H):
    c = S.STAGE_COLORS[i]
    return rrect(sl, x, y, w, h, fill=tint(c, 0.93), line=tint(c, 0.62), radius=0.09, lw=0.75)


STYLE = {
    'lead':   dict(size=11,  bold=True,  color=S.NAVY,    after=4.0),
    'norm':   dict(size=10,  bold=False, color=S.GREY,    after=4.0),
    'strong': dict(size=10,  bold=True,  color=None,      after=4.0),
    'mute':   dict(size=9,   bold=False, color=S.GREY_LT, after=3.0),
    'label':  dict(size=9,   bold=True,  color=None,      after=6.0, spc=50, caps=True),
}


def body_text(sl, i, x, w, lines, y=None, h=None, scale=1.0, anchor=MSO_ANCHOR.TOP):
    c = S.STAGE_COLORS[i]
    y = BOD_Y + 0.16 if y is None else y
    h = BOD_H - 0.32 if h is None else h
    paras = []
    for text, sty in lines:
        d = dict(STYLE[sty])
        if d['color'] is None:
            d['color'] = shade(c, 0.32)
        d['size'] = round(d['size'] * scale, 1)
        d['ls'] = 0.94
        d['text'] = text.upper() if d.pop('caps', False) else text
        paras.append(d)
    return tbox(sl, x + 0.15, y, w - 0.30, h, paras, anchor=anchor)


def add_chips(sl, chips, y=CHIP_Y, h=CHIP_H):
    out, g = [], 0.24
    w = (CW - 2 * g) / 3.0
    for j, txt in enumerate(chips):
        x = ML + j * (w + g)
        out.append(rrect(sl, x, y, w, h, fill=tint(S.NAVY, 0.93), radius=0.09))
        out.append(oval(sl, x + 0.22, y + h / 2 - 0.065, 0.13, S.GOLD))
        out.append(tbox(sl, x + 0.47, y + 0.03, w - 0.62, h - 0.06,
                        [{'text': txt, 'size': 10.5, 'color': S.NAVY, 'ls': 0.92}],
                        anchor=MSO_ANCHOR.MIDDLE))
    return out


def add_chevrons(sl):
    y = HDR_Y + HDR_H / 2 - 0.11
    out = []
    for k in range(len(COLX) - 1):
        gx = COLX[k] + COLW[k]
        gw = COLX[k + 1] - gx
        out.append(chevron(sl, gx + gw / 2 - 0.115, y, 0.23, 0.22, tint(S.NAVY, 0.70)))
    return out


def cpl_star(sl):
    return star(sl, COLX[2] + COLW[2] - 0.38, HDR_Y - 0.19, 0.42, S.GOLD)


# ------------------------------------------------------------------ spine
def credit_ladder(sl, y):
    """Stacked 120-unit bar: what CPL is worth against the two degrees."""
    out = []
    out.append(tbox(sl, ML, y, CW, 0.24,
                    [{'text': 'WHAT THE CREDIT ADDS UP TO', 'size': 9.5, 'bold': True,
                      'color': S.NAVY, 'spc': 60}]))
    by, bh = y + 0.32, 0.54
    segs = [(32.0, S.S3, 'CPL AWARD', '26–38 units'),
            (28.0, S.S4, 'REST OF THE ASSOCIATE DEGREE', 'balance of the 60'),
            (60.0, S.S5, "UPPER-DIVISION BACHELOR'S COURSEWORK", '60 units')]
    x = ML
    for units, color, name, detail in segs:
        w = CW * units / 120.0
        out.append(rect(sl, x, by, w, bh, color))
        out.append(tbox(sl, x + 0.14, by + 0.04, w - 0.28, bh - 0.08,
                        [{'text': name, 'size': 8.5, 'bold': True, 'color': S.WHITE,
                          'spc': 20, 'ls': 0.9, 'after': 1},
                         {'text': detail, 'size': 8.5, 'color': tint(color, 0.75), 'ls': 0.9}],
                        anchor=MSO_ANCHOR.MIDDLE))
        x += w
    ty = by + bh + 0.07
    b60 = ML + CW * 60.0 / 120.0
    out.append(tbox(sl, ML, ty, 1.2, 0.20,
                    [{'text': '0 units', 'size': 8.5, 'color': S.GREY_LT}]))
    out.append(tbox(sl, b60 - 1.6, ty, 3.2, 0.20,
                    [{'text': '60 units  ·  ASSOCIATE DEGREE', 'size': 8.5, 'bold': True,
                      'color': S.GREY, 'align': PP_ALIGN.CENTER, 'spc': 20}]))
    out.append(tbox(sl, ML + CW - 3.4, ty, 3.4, 0.20,
                    [{'text': "120 units  ·  BACHELOR'S DEGREE", 'size': 8.5, 'bold': True,
                      'color': S.GREY, 'align': PP_ALIGN.RIGHT, 'spc': 20}]))
    return out


def build_spine(prs, layout):
    sl = new_slide(prs, layout)
    add_heading(sl, S.SPINE['kicker'], S.SPINE['title'], S.SPINE['sub'], title_size=30)

    L1, L2 = [], []
    ch = 1.50
    for i in range(5):
        x, w = COLX[i], COLW[i]
        L1 += stage_header(sl, i, x, w)
        L2.append(body_card(sl, i, x, w, y=BOD_Y, h=ch))
        L2.append(tbox(sl, x + 0.16, BOD_Y + 0.16, w - 0.32, ch - 0.32,
                       [{'text': S.SPINE['stages'][i], 'size': 10.5, 'color': S.GREY,
                         'ls': 1.0}]))
    L1 += add_chevrons(sl)
    L1.append(cpl_star(sl))

    L3 = credit_ladder(sl, BOD_Y + ch + 0.30)
    L3 += add_chips(sl, S.SPINE['chips'])
    return sl, [L1, L2, L3]


# ------------------------------------------------------------------ trade
def build_trade(prs, layout, T):
    sl = new_slide(prs, layout)
    add_heading(sl, T['kicker'], T['title'], T['sub'])

    L1, L2, L3 = [], [], []
    for i in range(5):
        L1 += stage_header(sl, i, COLX[i], COLW[i])
        L2.append(body_card(sl, i, COLX[i], COLW[i]))
    L1 += add_chevrons(sl)

    for i in (0, 1, 3, 4):
        L2.append(body_text(sl, i, COLX[i], COLW[i], T['bodies'][i]))

    if T.get('logo'):
        L2.append(sl.shapes.add_picture(
            T['logo'], Inches(COLX[1] + COLW[1] / 2 - 0.47),
            Inches(BOD_Y + BOD_H - 1.06), height=Inches(0.94)))

    cx, cw = COLX[2], COLW[2]
    num, lab, note = T['hero']
    L3.append(cpl_star(sl))
    L3.append(tbox(sl, cx + 0.16, BOD_Y + 0.08, 1.16, 0.80,
                   [{'text': num, 'size': 46, 'bold': True,
                     'color': shade(S.S3, 0.22), 'ls': 0.85}]))
    L3.append(tbox(sl, cx + 1.34, BOD_Y + 0.18, cw - 1.50, 0.60,
                   [{'text': lab, 'size': 9, 'bold': True, 'color': S.GREY,
                     'ls': 0.95, 'spc': 20}], anchor=MSO_ANCHOR.MIDDLE))
    L3.append(tbox(sl, cx + 0.16, BOD_Y + 0.88, cw - 0.32, 0.22,
                   [{'text': note, 'size': 9, 'italic': True, 'color': S.GREY_LT}]))
    L2.append(rect(sl, cx + 0.16, BOD_Y + 1.16, cw - 0.32, 0.011, tint(S.S3, 0.55)))
    L2.append(body_text(sl, 2, cx, cw, T['bodies'][2],
                        y=BOD_Y + 1.30, h=BOD_H - 1.46))

    L3 += add_chips(sl, T['chips'])
    return sl, [L1, L2, L3]


# ------------------------------------------------------------------ ironworkers
def build_ironworkers(prs, layout, T):
    sl = new_slide(prs, layout)
    add_heading(sl, T['kicker'], T['title'], T['sub'])

    L1, L2, L3 = [], [], []
    for i in range(5):
        L1 += stage_header(sl, i, COLX[i], COLW[i])
    L1 += add_chevrons(sl)
    L1.append(cpl_star(sl))

    L2.append(body_card(sl, 0, COLX[0], COLW[0]))
    L2.append(body_text(sl, 0, COLX[0], COLW[0], T['stage1']))
    if T.get('logo'):
        L2.append(sl.shapes.add_picture(
            T['logo'], Inches(COLX[0] + COLW[0] / 2 - 0.50),
            Inches(BOD_Y + BOD_H - 1.10), height=Inches(0.98)))

    for li, lane in enumerate(T['lanes']):
        ly = LANE_A_Y if li == 0 else LANE_B_Y
        for k, cell in enumerate(lane['cells']):
            i = k + 1
            x, w = COLX[i], COLW[i]
            L2.append(body_card(sl, i, x, w, y=ly, h=LANE_H))
            if i == 2:
                L3.append(tbox(sl, x + 0.16, ly + 0.09, w - 0.32, 0.40,
                               [{'runs': [
                                   {'text': lane['units'], 'size': 26, 'bold': True,
                                    'color': shade(S.S3, 0.22)},
                                   {'text': '  UNITS', 'size': 10, 'bold': True,
                                    'color': S.GREY, 'spc': 40}], 'ls': 0.9}],
                               anchor=MSO_ANCHOR.MIDDLE))
                L2.append(body_text(sl, i, x, w, cell[1:], y=ly + 0.54,
                                    h=LANE_H - 0.66, scale=0.90))
            else:
                L2.append(body_text(sl, i, x, w, cell, y=ly + 0.10,
                                    h=LANE_H - 0.20, scale=0.92,
                                    anchor=MSO_ANCHOR.MIDDLE))

    gy = LANE_A_Y + LANE_H + 0.02
    ox = COLX[1] + COLW[1] / 2 - 0.33
    L1.append(rrect(sl, ox, gy, 0.66, 0.34, fill=S.NAVY, radius=0.17))
    L1.append(tbox(sl, ox, gy, 0.66, 0.34,
                   [{'text': 'OR', 'size': 10.5, 'bold': True, 'color': S.WHITE,
                     'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE))

    L3 += add_chips(sl, T['chips'])
    return sl, [L1, L2, L3]


# ------------------------------------------------------------------ appendix
def course_table(sl, x, y, w, rows, headers, col_w, fs=8.5, rh=0.235, hdr=S.S3):
    gf = sl.shapes.add_table(len(rows) + 1, 3, Inches(x), Inches(y), Inches(w), Inches(0.001))
    tbl = gf.table
    tbl.first_row = True
    for j, cwd in enumerate(col_w):
        tbl.columns[j].width = Inches(cwd)
    for ri, row in enumerate([list(headers)] + [list(r) for r in rows]):
        tbl.rows[ri].height = Inches(rh if ri else rh + 0.02)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                hdr if ri == 0 else (S.WHITE if ri % 2 else tint(hdr, 0.94)))
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if ci == 2 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fs)
            r.font.name = S.BODY_FONT
            r.font.bold = (ri == 0) or (ci == 0 and ri > 0)
            r.font.color.rgb = RGBColor.from_string(
                S.WHITE if ri == 0 else (S.NAVY if ci == 0 else S.GREY))
    return gf


def note(sl, x, y, w, text, size=8.5):
    return tbox(sl, x, y, w, 0.42,
                [{'text': text, 'size': size, 'italic': True,
                  'color': S.GREY_LT, 'ls': 1.0}])


def build_appendix_carp(prs, layout):
    sl = new_slide(prs, layout)
    add_heading(sl, 'APPENDIX  ·  CPL CROSSWALK', 'Carpentry — Drywall / Lather',
                'Santiago Canyon College course equivalents for Western States Carpenters '
                'apprenticeship training.', title_size=25)
    course_table(sl, ML, 1.52, 7.55, S.CARP_COURSES,
                 ('Course', 'Title', 'Units'), (1.35, 5.10, 1.10), rh=0.245)
    note(sl, ML, 4.84, 7.55,
         'As listed on the source pathway diagram: 12 courses, 20.0 units.')

    px, pw = 8.35, 4.58
    rrect(sl, px, 1.52, pw, 3.16, fill=tint(S.NAVY, 0.94), radius=0.10)
    tbox(sl, px + 0.28, 1.76, pw - 0.56, 2.70,
         [{'text': 'HOW TO READ THIS', 'size': 9.5, 'bold': True, 'color': S.NAVY,
           'spc': 60, 'after': 9},
          {'text': 'Each row is training the apprentice has already completed, matched to a '
                   'college course the college has already approved.', 'size': 10.5,
           'color': S.GREY, 'ls': 1.02, 'after': 9},
          {'text': 'Nothing here is new curriculum. The crosswalk is what turns completed '
                   'apprenticeship training into transcripted college credit.', 'size': 10.5,
           'color': S.GREY, 'ls': 1.02, 'after': 9},
          {'text': 'The pathway slide cites 26 units eligible at journeyperson. The courses '
                   'making up the difference ran off the bottom edge of the original graphic '
                   '— confirm the full list before this page is shared.', 'size': 10,
           'color': shade(S.S3, 0.30), 'ls': 1.02, 'italic': True}])
    return sl


def build_appendix_iron(prs, layout):
    sl = new_slide(prs, layout)
    add_heading(sl, 'APPENDIX  ·  CPL CROSSWALK', 'Ironworkers — two college routes',
                'The same trade training, evaluated by two colleges — Ironworkers Local 433.',
                title_size=25)
    tbox(sl, ML, 1.48, 6.0, 0.24,
         [{'text': 'AMERICAN RIVER COLLEGE   ·   29.5 UNITS', 'size': 10, 'bold': True,
           'color': S.NAVY, 'spc': 30}])
    course_table(sl, ML, 1.76, 6.05, S.ARC_COURSES,
                 ('ARC', 'Ironworker course', 'Units'), (0.85, 4.35, 0.85), fs=7.5, rh=0.225)
    tbox(sl, 6.93, 1.48, 6.0, 0.24,
         [{'text': 'CERRITOS COLLEGE   ·   38 UNITS', 'size': 10, 'bold': True,
           'color': S.NAVY, 'spc': 30}])
    course_table(sl, 6.93, 1.76, 6.0, S.CC_COURSES,
                 ('CC', 'Course title', 'Units'), (1.30, 3.85, 0.85), fs=7.5, rh=0.225)
    note(sl, 6.93, 5.60, 6.0,
         'Cerritos list as shown on the source diagram totals 31.5 units against the 38 cited; '
         'rows below the slide edge were not visible in the original graphic. Confirm before '
         'sharing. The American River list reconciles exactly at 29.5.')
    return sl


# ------------------------------------------------------------------ animation
def inject_timing(slide, layers, dur=400, stagger=55):
    ids = iter(range(3, 100000))
    groups = []
    for layer in layers:
        eff = []
        for k, shp in enumerate(layer):
            sid = shp.shape_id
            node = 'clickEffect' if k == 0 else 'withEffect'
            delay = 0 if k == 0 else min(k * stagger, 700)
            a, b, c = next(ids), next(ids), next(ids)
            eff.append(f'''
<p:par><p:cTn id="{a}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="{node}">
  <p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>
  <p:childTnLst>
    <p:set><p:cBhvr>
      <p:cTn id="{b}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
      <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
    </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>
    <p:animEffect transition="in" filter="fade">
      <p:cBhvr><p:cTn id="{c}" dur="{dur}"/><p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl></p:cBhvr>
    </p:animEffect>
  </p:childTnLst></p:cTn></p:par>''')
        g1, g2 = next(ids), next(ids)
        groups.append(f'''
<p:par><p:cTn id="{g1}" fill="hold">
  <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
  <p:childTnLst>
    <p:par><p:cTn id="{g2}" fill="hold">
      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
      <p:childTnLst>{''.join(eff)}</p:childTnLst>
    </p:cTn></p:par>
  </p:childTnLst></p:cTn></p:par>''')

    xml = f'''<p:timing xmlns:p="{NS['p']}" xmlns:a="{NS['a']}">
<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
{''.join(groups)}
</p:childTnLst></p:cTn>
<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'''

    el = etree.fromstring(xml)
    sld = slide._element
    for old in sld.findall(qn('p:timing')):
        sld.remove(old)
    ext = sld.find(qn('p:extLst'))
    if ext is not None:
        ext.addprevious(el)
    else:
        sld.append(el)


# ------------------------------------------------------------------ assembly
def reorder(prs, order):
    lst = prs.slides._sldIdLst
    items = list(lst)
    for it in items:
        lst.remove(it)
    for idx in order:
        lst.append(items[idx])


def main(animate=True, out='out.pptx'):
    prs = Presentation('source.pptx')
    layout = prs.slide_masters[0].slide_layouts[2]     # 'Title and Content (No Symbol)'

    built = [build_spine(prs, layout),
             build_trade(prs, layout, S.CARPENTRY),
             build_ironworkers(prs, layout, S.IRONWORKERS),
             build_trade(prs, layout, S.FIRE)]
    build_appendix_carp(prs, layout)
    build_appendix_iron(prs, layout)

    if animate:
        for sl, layers in built:
            inject_timing(sl, layers)

    # keep the three originals in the file, hidden at the end -- nothing is discarded
    order = [0, 1, 2, 3, 4, 12, 13, 14, 15, 8, 9, 10, 11, 16, 17, 5, 6, 7]
    reorder(prs, order)
    for idx in (5, 6, 7):
        prs.slides[order.index(idx)]._element.set('show', '0')
    prs.save(out)
    shown = sum(1 for s in prs.slides if s._element.get('show') != '0')
    print('wrote', out, '| slides =', len(prs.slides._sldIdLst), '| shown =', shown)


if __name__ == '__main__':
    main(animate=('--static' not in sys.argv),
         out=sys.argv[sys.argv.index('-o') + 1] if '-o' in sys.argv else 'out.pptx')
